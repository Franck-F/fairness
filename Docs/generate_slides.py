# -*- coding: utf-8 -*-
"""
Generation des slides de soutenance Epitech Digital School pour AuditIQ.
Utilise python-pptx. Sortie : Docs/soutenance_slides.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

import os

# =====================================================================
# Palette
# =====================================================================
NAVY = RGBColor(0x1E, 0x3A, 0x8A)      # bleu nuit
ORANGE = RGBColor(0xF5, 0x9E, 0x0B)    # accent
LIGHT_BG = RGBColor(0xF8, 0xFA, 0xFC)  # gris tres clair
DARK_BG = RGBColor(0x0F, 0x17, 0x2A)   # fond sombre
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_MID = RGBColor(0x64, 0x74, 0x8B)
GRAY_DARK = RGBColor(0x1F, 0x29, 0x37)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

# =====================================================================
# Presentation 16:9
# =====================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

BLANK = prs.slide_layouts[6]


def add_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg


def add_text(slide, left, top, width, height, text, *,
             size=18, bold=False, color=GRAY_DARK, font=FONT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return tb


def add_rect(slide, left, top, width, height, fill=None, line=None, line_w=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_round_rect(slide, left, top, width, height, fill=None, line=None, line_w=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.adjustments[0] = 0.12
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_oval(slide, left, top, width, height, fill=None, line=None, line_w=0, alpha=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
        if alpha is not None:
            # set alpha on the solidFill element
            sp = shp.fill.fore_color._xFill
            srgb = sp.find(qn('a:srgbClr'))
            if srgb is not None:
                alpha_el = srgb.makeelement(qn('a:alpha'),
                                            {'val': str(int(alpha * 100000))})
                srgb.append(alpha_el)
    else:
        shp.fill.background()
    if line is not None:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_footer(slide, dark=False):
    col = WHITE if dark else GRAY_MID
    add_text(slide, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.3),
             "Franck F. — Epitech Digital School 2024-2025",
             size=10, color=col)


# =====================================================================
# SLIDE 1 — TITRE (fond sombre)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, DARK_BG)

# Barre orange a gauche
add_rect(s, 0, Inches(2.2), Inches(0.35), Inches(3.2), fill=ORANGE)

add_text(s, Inches(0.8), Inches(0.8), Inches(12), Inches(0.5),
         "CONSULTING PROJECT — MEMOIRE DE FIN D'ETUDES",
         size=14, bold=True, color=ORANGE)

add_text(s, Inches(0.8), Inches(1.5), Inches(12), Inches(1.3),
         "AuditIQ",
         size=72, bold=True, color=WHITE)

add_text(s, Inches(0.8), Inches(2.7), Inches(12), Inches(1.6),
         "Rendre l'audit de fairness algorithmique\naccessible aux PME francaises",
         size=30, bold=False, color=WHITE)

# Sous-titre question
add_text(s, Inches(0.8), Inches(4.6), Inches(12), Inches(1.3),
         "Comment permettre aux PME francaises de detecter facilement\n"
         "les biais de leurs algorithmes IA pour se conformer\n"
         "a l'AI Act europeen ?",
         size=18, color=RGBColor(0xCB, 0xD5, 0xE1))

# Auteur
add_rect(s, Inches(0.8), Inches(6.4), Inches(5), Inches(0.03), fill=ORANGE)
add_text(s, Inches(0.8), Inches(6.5), Inches(8), Inches(0.4),
         "Franck F.  •  MSc 2024-2025  •  Epitech Digital School",
         size=14, bold=True, color=WHITE)

add_footer(s, dark=True)


# =====================================================================
# SLIDE 2 — LE PROBLEME EN CHIFFRES (fond clair)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, LIGHT_BG)

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "01  —  LE CONSTAT", size=12, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(0.9), Inches(12.5), Inches(1),
         "Le probleme en chiffres",
         size=40, bold=True, color=NAVY)

# 4 cartes
cards = [
    ("38 %", "des PME ignorent\nce qu'est un biais\nalgorithmique"),
    ("44 %", "soupconnent que\nleurs outils IA peuvent\ndiscriminer"),
    ("35 %", "n'ont aucun budget\n(0 €) pour l'audit\nethique IA"),
    ("82 %", "disposent de moins\nde 2 000 € de budget\naudit IA"),
]
card_w = Inches(2.85)
card_h = Inches(3.8)
gap = Inches(0.2)
start_left = Inches(0.6)
top = Inches(2.3)

for i, (num, label) in enumerate(cards):
    left = start_left + (card_w + gap) * i
    add_round_rect(s, left, top, card_w, card_h, fill=WHITE)
    # Barre superieure orange
    add_rect(s, left + Inches(0.3), top + Inches(0.35),
             Inches(0.7), Inches(0.08), fill=ORANGE)
    add_text(s, left, top + Inches(0.7), card_w, Inches(1.6),
             num, size=64, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, left + Inches(0.2), top + Inches(2.4),
             card_w - Inches(0.4), Inches(1.3),
             label, size=15, color=GRAY_DARK, align=PP_ALIGN.CENTER)

# Source
add_text(s, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "Source : Sondage « IA et Ethique dans les PME » — 34 repondants (2025)",
         size=11, color=GRAY_MID)

add_footer(s)


# =====================================================================
# SLIDE 3 — LE TRIPLE FOSSE (fond clair)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, LIGHT_BG)

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "02  —  LE PARADOXE", size=12, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(0.9), Inches(12.5), Inches(1),
         "Le triple fosse",
         size=40, bold=True, color=NAVY)

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
         "Fairlearn existe. AIF360 existe. L'AI Act existe. Mais rien ne les connecte pour les PME.",
         size=16, color=GRAY_MID)

# Trois colonnes
cols = [
    ("01", "Fosse cognitif",
     "Les dirigeants PME ne\nconnaissent pas le vocabulaire :\nDemographic Parity, Equal\nOpportunity, disparate impact."),
    ("02", "Fosse technique",
     "Fairlearn et AIF360 sont des\nframeworks Python. Ils exigent\ndata scientist, notebooks,\nscikit-learn — hors de portee."),
    ("03", "Fosse reglementaire",
     "L'AI Act art. 10 impose\nl'examen des biais mais ne dit\npas comment. Aucune norme\noperationnelle pour les PME."),
]
col_w = Inches(4.0)
col_h = Inches(3.8)
gap = Inches(0.15)
start_left = Inches(0.6)
top = Inches(2.6)

for i, (num, title, body) in enumerate(cols):
    left = start_left + (col_w + gap) * i
    add_round_rect(s, left, top, col_w, col_h, fill=WHITE)
    add_text(s, left + Inches(0.3), top + Inches(0.3), Inches(1.2), Inches(0.9),
             num, size=44, bold=True, color=ORANGE)
    add_rect(s, left + Inches(0.3), top + Inches(1.25),
             Inches(0.8), Inches(0.05), fill=NAVY)
    add_text(s, left + Inches(0.3), top + Inches(1.45),
             col_w - Inches(0.6), Inches(0.6),
             title, size=22, bold=True, color=NAVY)
    add_text(s, left + Inches(0.3), top + Inches(2.15),
             col_w - Inches(0.6), Inches(2),
             body, size=14, color=GRAY_DARK)

# Conclusion bandeau
add_rect(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5), fill=NAVY)
add_text(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.5),
         "Les briques existent. Personne ne les connecte pour les PME.",
         size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_footer(s)


# =====================================================================
# SLIDE 4 — LA TRIPLE INTERFACE (fond sombre) — Venn
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, DARK_BG)

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "03  —  CONTRIBUTION THEORIQUE", size=12, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(0.9), Inches(12.5), Inches(1),
         "Le concept de triple interface",
         size=40, bold=True, color=WHITE)

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
         "Cognitive. Technique. Reglementaire. Une grille de lecture du fosse.",
         size=16, color=RGBColor(0xCB, 0xD5, 0xE1))

# Venn : 3 cercles sur la gauche
venn_left = Inches(0.8)
venn_top = Inches(2.4)
circle_d = Inches(3.0)

# Positionnement : triangle
c1_left = venn_left + Inches(0.0)
c1_top = venn_top + Inches(0.0)
c2_left = venn_left + Inches(1.6)
c2_top = venn_top + Inches(0.0)
c3_left = venn_left + Inches(0.8)
c3_top = venn_top + Inches(1.4)

add_oval(s, c1_left, c1_top, circle_d, circle_d,
         fill=RGBColor(0x3B, 0x82, 0xF6), alpha=0.55)
add_oval(s, c2_left, c2_top, circle_d, circle_d,
         fill=ORANGE, alpha=0.55)
add_oval(s, c3_left, c3_top, circle_d, circle_d,
         fill=RGBColor(0x10, 0xB9, 0x81), alpha=0.55)

# Labels
add_text(s, c1_left - Inches(0.4), c1_top - Inches(0.35),
         Inches(3), Inches(0.4),
         "COGNITIVE", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, c2_left + Inches(0.4), c2_top - Inches(0.35),
         Inches(3), Inches(0.4),
         "TECHNIQUE", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, c3_left - Inches(0.1), c3_top + circle_d + Inches(0.1),
         Inches(3.2), Inches(0.4),
         "REGLEMENTAIRE", size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Centre AuditIQ
add_text(s, venn_left + Inches(0.3), venn_top + Inches(1.6),
         Inches(3.8), Inches(0.6),
         "AuditIQ", size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Partie droite : explications
right_left = Inches(6.8)
add_text(s, right_left, Inches(2.5), Inches(6), Inches(0.5),
         "Une integration originale", size=20, bold=True, color=ORANGE)

rows = [
    ("Cognitive", "Langage naturel, feux tricolores, zero jargon."),
    ("Technique", "Fairlearn empaquete, wizard cinq etapes, upload CSV."),
    ("Reglementaire", "Mapping AI Act art. 10-11 + droit francais."),
]
ry = Inches(3.15)
for title, body in rows:
    add_rect(s, right_left, ry, Inches(0.08), Inches(0.7), fill=ORANGE)
    add_text(s, right_left + Inches(0.2), ry - Inches(0.05),
             Inches(6), Inches(0.35),
             title, size=14, bold=True, color=WHITE)
    add_text(s, right_left + Inches(0.2), ry + Inches(0.3),
             Inches(6), Inches(0.4),
             body, size=13, color=RGBColor(0xCB, 0xD5, 0xE1))
    ry += Inches(0.85)

# Citations
add_text(s, right_left, Inches(5.85), Inches(6), Inches(0.35),
         "Inspirations : Wachter, Mittelstadt & Russell (2021) ;",
         size=11, color=GRAY_MID)
add_text(s, right_left, Inches(6.05), Inches(6), Inches(0.35),
         "Bhatt et al. (2020) sur l'explicabilite contextuelle.",
         size=11, color=GRAY_MID)

# Phrase d'accroche bas
add_rect(s, Inches(0.6), Inches(6.7), Inches(12.1), Inches(0.5), fill=ORANGE)
add_text(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.5),
         "Empaqueter Fairlearn n'est pas une innovation algorithmique. Le faire avec une triple interface, c'est une innovation strategique.",
         size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

add_footer(s, dark=True)


# =====================================================================
# SLIDE 5 — MODULE 1 : AUDIT SUPERVISE (fond clair)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, LIGHT_BG)

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "04  —  MODULE 1", size=12, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(0.9), Inches(12.5), Inches(1),
         "Audit supervise classique",
         size=40, bold=True, color=NAVY)

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
         "Metriques de fairness standardisees sur un jeu de donnees labellise.",
         size=16, color=GRAY_MID)

# Zone placeholder capture d'ecran
add_round_rect(s, Inches(0.6), Inches(2.6), Inches(7.2), Inches(4.3),
               fill=WHITE, line=NAVY, line_w=1)
add_text(s, Inches(0.6), Inches(4.3), Inches(7.2), Inches(0.5),
         "[ CAPTURE D'ECRAN — wizard 5 etapes ]",
         size=16, bold=True, color=GRAY_MID, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(4.8), Inches(7.2), Inches(0.4),
         "a remplacer par Franck",
         size=11, color=GRAY_MID, align=PP_ALIGN.CENTER)

# Points cles
features = [
    ("Metriques", "Demographic Parity, Equal Opportunity, Equalized Odds, 4/5 rule."),
    ("Feux tricolores", "Lecture du risque en deux secondes, seuils pre-configures."),
    ("Langage naturel", "Explications generees pour un non-specialiste. Pas de jargon."),
]
fx = Inches(8.1)
fy = Inches(2.7)
for title, body in features:
    add_round_rect(s, fx, fy, Inches(4.9), Inches(1.35), fill=WHITE)
    add_rect(s, fx, fy, Inches(0.1), Inches(1.35), fill=ORANGE)
    add_text(s, fx + Inches(0.3), fy + Inches(0.15),
             Inches(4.6), Inches(0.4),
             title, size=15, bold=True, color=NAVY)
    add_text(s, fx + Inches(0.3), fy + Inches(0.55),
             Inches(4.6), Inches(0.75),
             body, size=12, color=GRAY_DARK)
    fy += Inches(1.5)

add_footer(s)


# =====================================================================
# SLIDE 6 — MODULE 2 : DETECTION NON SUPERVISEE (fond clair)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, LIGHT_BG)

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "05  —  MODULE 2  —  NOUVEAUTE", size=12, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(0.9), Inches(12.5), Inches(1),
         "Detection non supervisee",
         size=40, bold=True, color=NAVY)

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
         "Detecter des disparites sans collecter d'attributs sensibles.",
         size=16, color=GRAY_MID)

# Pipeline : 5 etapes
steps = ["CSV", "KMeans", "Khi²", "Caracterisation", "Verdict"]
step_w = Inches(2.2)
step_h = Inches(1.2)
step_gap = Inches(0.25)
total_w = step_w * 5 + step_gap * 4
start_x = (SW - total_w) / 2
py = Inches(3.1)

for i, label in enumerate(steps):
    x = start_x + (step_w + step_gap) * i
    add_round_rect(s, x, py, step_w, step_h, fill=WHITE, line=NAVY, line_w=1.5)
    add_text(s, x, py + Inches(0.35), step_w, step_h,
             label, size=18, bold=True, color=NAVY,
             align=PP_ALIGN.CENTER)
    if i < 4:
        arrow = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            x + step_w + Inches(0.02), py + Inches(0.45),
            Inches(0.22), Inches(0.3),
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ORANGE
        arrow.line.fill.background()

# Inspiration
add_round_rect(s, Inches(0.6), Inches(4.9), Inches(6), Inches(1.9),
               fill=WHITE)
add_rect(s, Inches(0.6), Inches(4.9), Inches(0.1), Inches(1.9), fill=NAVY)
add_text(s, Inches(0.85), Inches(5.0), Inches(5.7), Inches(0.4),
         "INSPIRATION", size=11, bold=True, color=ORANGE)
add_text(s, Inches(0.85), Inches(5.3), Inches(5.7), Inches(0.5),
         "Algorithm Audit (Pays-Bas)",
         size=16, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(5.75), Inches(5.7), Inches(1),
         "Approche reconnue par l'OCDE\nCatalogue of Tools for Trustworthy AI.",
         size=12, color=GRAY_DARK)

# Argument cle
add_round_rect(s, Inches(6.9), Inches(4.9), Inches(5.8), Inches(1.9),
               fill=NAVY)
add_text(s, Inches(7.1), Inches(5.0), Inches(5.5), Inches(0.4),
         "ARGUMENT CLE", size=11, bold=True, color=ORANGE)
add_text(s, Inches(7.1), Inches(5.3), Inches(5.5), Inches(1.5),
         "Resout le paradoxe AI Act art. 10(5)\n"
         "contre RGPD : pas besoin de collecter\n"
         "d'attributs sensibles.",
         size=13, bold=True, color=WHITE)

add_footer(s)


# =====================================================================
# SLIDE 7 — MODULE 3 : CHATBOT/LLM (fond clair)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, LIGHT_BG)

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "06  —  MODULE 3", size=12, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(0.9), Inches(12.5), Inches(1),
         "Audit chatbot / LLM",
         size=40, bold=True, color=NAVY)

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
         "17 % des PME du sondage utilisent un chatbot. Aucune n'a audite ses biais.",
         size=16, color=GRAY_MID)

# Grand chiffre
add_round_rect(s, Inches(0.6), Inches(2.7), Inches(4), Inches(4.2),
               fill=NAVY)
add_text(s, Inches(0.6), Inches(3.0), Inches(4), Inches(1.5),
         "17 %", size=96, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.0), Inches(3.6), Inches(1.7),
         "des PME du panel\nutilisent un assistant\nconversationnel en production",
         size=14, color=WHITE, align=PP_ALIGN.CENTER)

# Specs a droite
specs_left = Inches(4.9)
add_round_rect(s, specs_left, Inches(2.7), Inches(7.8), Inches(4.2),
               fill=WHITE)
add_text(s, specs_left + Inches(0.35), Inches(2.85),
         Inches(7), Inches(0.4),
         "INSPIRATION", size=11, bold=True, color=ORANGE)
add_text(s, specs_left + Inches(0.35), Inches(3.15),
         Inches(7), Inches(0.5),
         "LangBiTe — UOC Barcelone / Universite du Luxembourg",
         size=15, bold=True, color=NAVY)

add_rect(s, specs_left + Inches(0.35), Inches(3.7),
         Inches(7), Inches(0.03), fill=ORANGE)

specs = [
    ("10", "prompts paired bilingues (FR / EN)"),
    ("06", "categories de biais (genre, age, origine, religion, handicap, orientation)"),
    ("HTTP", "adaptateur generique compatible OpenAI, Anthropic, Mistral"),
]
sy = Inches(3.9)
for num, label in specs:
    add_text(s, specs_left + Inches(0.35), sy, Inches(1.3), Inches(0.7),
             num, size=26, bold=True, color=ORANGE)
    add_text(s, specs_left + Inches(1.7), sy + Inches(0.15),
             Inches(6), Inches(0.5),
             label, size=13, color=GRAY_DARK)
    sy += Inches(0.95)

add_footer(s)


# =====================================================================
# SLIDE 8 — DEMO (fond clair, minimaliste)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, LIGHT_BG)

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "07  —  DEMO", size=12, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(2.4), Inches(12.5), Inches(1.5),
         "Demonstration en direct",
         size=60, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

add_rect(s, Inches(5.5), Inches(3.8), Inches(2.3), Inches(0.06), fill=ORANGE)

add_text(s, Inches(0.6), Inches(4.1), Inches(12.5), Inches(0.6),
         "4 minutes",
         size=28, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.6), Inches(5.0), Inches(12.5), Inches(0.6),
         "Scenario : audit du dataset 01_recrutement_ats.csv",
         size=18, color=GRAY_DARK, align=PP_ALIGN.CENTER)

add_footer(s)


# =====================================================================
# SLIDE 9 — EVALUATION UTILISATEUR (fond clair)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, LIGHT_BG)

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "08  —  EVALUATION", size=12, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(0.9), Inches(12.5), Inches(1),
         "Protocole d'evaluation utilisateur",
         size=40, bold=True, color=NAVY)

add_text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(0.5),
         "Un protocole rigoureux, des chiffres mesures — pas affirmes.",
         size=16, color=GRAY_MID)

# 4 cartes methode
items = [
    ("SUS", "System Usability\nScale (Brooke 1996)"),
    ("Time-on-task", "Duree jusqu'au premier\nrapport d'audit genere"),
    ("Comparaison", "AuditIQ UI vs\nFairlearn CLI"),
    ("n = 6-8", "Echantillon Nielsen\n(testeurs PME)"),
]
iw = Inches(2.95)
ih = Inches(1.9)
gap = Inches(0.12)
ileft = Inches(0.6)
itop = Inches(2.7)
for i, (title, body) in enumerate(items):
    x = ileft + (iw + gap) * i
    add_round_rect(s, x, itop, iw, ih, fill=WHITE)
    add_rect(s, x, itop, iw, Inches(0.08), fill=ORANGE)
    add_text(s, x + Inches(0.2), itop + Inches(0.25),
             iw - Inches(0.4), Inches(0.5),
             title, size=18, bold=True, color=NAVY)
    add_text(s, x + Inches(0.2), itop + Inches(0.8),
             iw - Inches(0.4), Inches(1),
             body, size=12, color=GRAY_DARK)

# Zone placeholder resultats
add_round_rect(s, Inches(0.6), Inches(4.9), Inches(7.5), Inches(1.8),
               fill=WHITE, line=NAVY, line_w=1)
add_text(s, Inches(0.85), Inches(5.0), Inches(7), Inches(0.4),
         "RESULTATS MESURES", size=11, bold=True, color=ORANGE)
add_text(s, Inches(0.85), Inches(5.3), Inches(7), Inches(0.4),
         "SUS moyen : [ __ / 100 ]   •   Time-on-task median : [ __ min ]",
         size=14, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(5.7), Inches(7), Inches(0.4),
         "Completion : [ __ % ]   •   vs Fairlearn CLI : [ __ min ]",
         size=14, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(6.15), Inches(7), Inches(0.4),
         "a completer apres les sessions",
         size=10, color=GRAY_MID)

# Zone honnetete
add_round_rect(s, Inches(8.3), Inches(4.9), Inches(4.4), Inches(1.8),
               fill=NAVY)
add_text(s, Inches(8.5), Inches(5.0), Inches(4), Inches(0.4),
         "LUCIDITE METHODOLOGIQUE", size=10, bold=True, color=ORANGE)
add_text(s, Inches(8.5), Inches(5.3), Inches(4), Inches(1.4),
         "n = 34 modeste.\n"
         "Marges d'erreur +/- 17 pts a 95 %.\n"
         "Triangulation HBR 2025 et\n"
         "Accountancy Europe 2025.",
         size=11, color=WHITE)

add_footer(s)


# =====================================================================
# SLIDE 10 — CONCLUSION (fond sombre)
# =====================================================================
s = prs.slides.add_slide(BLANK)
add_bg(s, DARK_BG)

add_text(s, Inches(0.6), Inches(0.5), Inches(12), Inches(0.4),
         "09  —  CONCLUSION", size=12, bold=True, color=ORANGE)

add_text(s, Inches(0.6), Inches(0.9), Inches(12.5), Inches(1),
         "Contributions, limites, perspectives",
         size=36, bold=True, color=WHITE)

# 3 contributions
contribs = [
    ("Conceptuelle",
     "Le cadre de la triple\ninterface comme grille\nd'analyse."),
    ("Methodologique",
     "Combiner supervise,\nnon supervise et LLM\ndans une plateforme\nnon-experts."),
    ("Empirique",
     "Validation qualitative\naupres de 34 PME\nfrancophones."),
]
cw = Inches(4.0)
ch = Inches(2.4)
cgap = Inches(0.15)
cleft = Inches(0.6)
ctop = Inches(2.3)
for i, (title, body) in enumerate(contribs):
    x = cleft + (cw + cgap) * i
    add_round_rect(s, x, ctop, cw, ch,
                   fill=RGBColor(0x1E, 0x29, 0x3B))
    add_text(s, x + Inches(0.3), ctop + Inches(0.25),
             cw - Inches(0.6), Inches(0.4),
             f"0{i+1}", size=12, bold=True, color=ORANGE)
    add_text(s, x + Inches(0.3), ctop + Inches(0.55),
             cw - Inches(0.6), Inches(0.5),
             title, size=20, bold=True, color=WHITE)
    add_text(s, x + Inches(0.3), ctop + Inches(1.05),
             cw - Inches(0.6), Inches(1.3),
             body, size=13, color=RGBColor(0xCB, 0xD5, 0xE1))

# Limites et pistes
add_text(s, Inches(0.6), Inches(5.0), Inches(6), Inches(0.4),
         "LIMITES ASSUMEES", size=11, bold=True, color=ORANGE)
add_text(s, Inches(0.6), Inches(5.3), Inches(6), Inches(1.2),
         "n = 34, biais de recrutement,\n"
         "couverture AI Act art. 10-11\n"
         "(soit ~15 % des obligations).",
         size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

add_text(s, Inches(7), Inches(5.0), Inches(6), Inches(0.4),
         "PISTES THESE", size=11, bold=True, color=ORANGE)
add_text(s, Inches(7), Inches(5.3), Inches(6), Inches(1.2),
         "Articulation droit / fairness\n"
         "computationnelle, efficacite\n"
         "longitudinale, gouvernance\n"
         "algorithmique europeenne.",
         size=12, color=RGBColor(0xCB, 0xD5, 0xE1))

# Phrase finale bandeau
add_rect(s, Inches(0.6), Inches(6.75), Inches(12.1), Inches(0.45), fill=ORANGE)
add_text(s, Inches(0.6), Inches(6.78), Inches(12.1), Inches(0.4),
         "AuditIQ n'est pas un certificat de conformite. C'est un point d'entree documente. C'est ce qui manque aujourd'hui aux PME francaises.",
         size=12, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

add_footer(s, dark=True)


# =====================================================================
# Sauvegarde
# =====================================================================
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        "soutenance_slides.pptx")
prs.save(out_path)
print("OK:", out_path)
print("Slides:", len(prs.slides))
